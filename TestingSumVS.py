import os, json, re, time, requests
from io import BytesIO
from typing import List, Dict, Any, Tuple
import streamlit as st
from pptx import Presentation
from dotenv import load_dotenv
from openai import AzureOpenAI
import PyPDF2
from pptx.dml.color import RGBColor

# ========== Config ==========
load_dotenv()
AZURE_API_KEY  = os.getenv("AZURE_OPENAI_API_KEY")
AZURE_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
API_VERSION    = os.getenv("AZURE_OPENAI_API_VERSION", "2024-02-15-preview")
CHAT_MODEL     = os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME", "gpt-4o-mini")
DEBUG_MODE     = os.getenv("DEBUG_MODE", "False").lower() == "true"

# ========== Modern UI Styling ==========
def inject_custom_css():
    st.markdown("""
    <style>
    /* Import Google Fonts */
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');
    
    /* Global Styles */
    * {
        font-family: 'Inter', sans-serif;
    }
    
    /* Main container styling */
    .main {
        padding: 2rem 1rem;
    }
    
    /* Header styling */
    h1 {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        font-size: 3rem;
        font-weight: 700;
        margin-bottom: 0.5rem;
    }
    
    h2 {
        color: #1a202c;
        font-weight: 600;
        margin-top: 2rem;
        margin-bottom: 1rem;
    }
    
    h3 {
        color: #2d3748;
        font-weight: 600;
        margin-top: 1.5rem;
        margin-bottom: 0.75rem;
    }
    
    /* Card styling */
    .deck-card {
        background: white;
        border-radius: 16px;
        padding: 2rem;
        box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1), 0 2px 4px -1px rgba(0, 0, 0, 0.06);
        margin-bottom: 2rem;
        border: 1px solid #e2e8f0;
        transition: all 0.3s ease;
    }
    
    .deck-card:hover {
        box-shadow: 0 10px 15px -3px rgba(0, 0, 0, 0.1), 0 4px 6px -2px rgba(0, 0, 0, 0.05);
        transform: translateY(-2px);
    }
    
    /* Metrics card */
    .metric-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        border-radius: 12px;
        padding: 1.5rem;
        color: white;
        text-align: center;
        box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1);
    }
    
    .metric-value {
        font-size: 2.5rem;
        font-weight: 700;
        margin-bottom: 0.25rem;
    }
    
    .metric-label {
        font-size: 0.875rem;
        opacity: 0.9;
        text-transform: uppercase;
        letter-spacing: 0.05em;
    }
    
    /* Context badges */
    .context-badge {
        background: linear-gradient(135deg, #48bb78 0%, #38a169 100%);
        color: white;
        border-radius: 24px;
        padding: 0.5rem 1.25rem;
        font-size: 0.875rem;
        font-weight: 600;
        display: inline-block;
        margin: 0.25rem;
        box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
        text-transform: capitalize;
    }
    
    .context-badge.purple {
        background: linear-gradient(135deg, #9f7aea 0%, #805ad5 100%);
    }
    
    .context-badge.blue {
        background: linear-gradient(135deg, #4299e1 0%, #3182ce 100%);
    }
    
    /* Summary box */
    .summary-box {
        background: #f7fafc;
        border-left: 4px solid #667eea;
        padding: 1.5rem;
        border-radius: 8px;
        margin: 1rem 0;
        line-height: 1.8;
    }
    
    /* AI suggestion box */
    .ai-suggestion {
        background: linear-gradient(135deg, #f6f9fc 0%, #e9f5ff 100%);
        border: 1px solid #bee3f8;
        padding: 1.25rem;
        border-radius: 12px;
        margin: 1rem 0;
        position: relative;
    }
    
    .ai-suggestion::before {
        content: "✨ AI Suggestion";
        position: absolute;
        top: -10px;
        left: 20px;
        background: white;
        padding: 0 8px;
        font-size: 0.75rem;
        color: #3182ce;
        font-weight: 600;
    }
    
    /* Color swatch */
    .color-palette {
        display: flex;
        gap: 0.5rem;
        margin: 1rem 0;
        padding: 1rem;
        background: #f7fafc;
        border-radius: 8px;
        align-items: center;
        flex-wrap: wrap;
    }
    
    .color-swatch {
        width: 60px;
        height: 60px;
        border-radius: 8px;
        box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
        border: 2px solid white;
        position: relative;
        cursor: pointer;
        transition: transform 0.2s ease;
    }
    
    .color-swatch:hover {
        transform: scale(1.1);
    }
    
    .color-hex {
        position: absolute;
        bottom: -20px;
        left: 50%;
        transform: translateX(-50%);
        font-size: 0.7rem;
        color: #4a5568;
        font-weight: 500;
    }
    
    /* Progress bar */
    .stProgress > div > div > div > div {
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
    }
    
    /* Buttons */
    .stButton > button {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        border: none;
        padding: 0.75rem 2rem;
        font-weight: 600;
        border-radius: 8px;
        transition: all 0.3s ease;
        box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1);
    }
    
    .stButton > button:hover {
        box-shadow: 0 10px 15px -3px rgba(0, 0, 0, 0.1);
        transform: translateY(-2px);
    }
    
    /* Sidebar styling */
    .css-1d391kg {
        background: #f7fafc;
    }
    
    /* Tab styling */
    .stTabs [data-baseweb="tab-list"] {
        gap: 2rem;
        background-color: #f7fafc;
        padding: 0.5rem;
        border-radius: 12px;
    }
    
    .stTabs [data-baseweb="tab"] {
        height: 3rem;
        white-space: pre-wrap;
        background-color: transparent;
        border-radius: 8px;
        padding: 0 1.5rem;
        font-weight: 600;
        color: #4a5568;
        transition: all 0.3s ease;
    }
    
    .stTabs [data-baseweb="tab"]:hover {
        color: #667eea;
        background-color: rgba(102, 126, 234, 0.1);
    }
    
    .stTabs [aria-selected="true"] {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white !important;
    }
    
    /* File uploader */
    .stFileUploader {
        background: white;
        border: 2px dashed #cbd5e0;
        border-radius: 12px;
        padding: 2rem;
        text-align: center;
        transition: all 0.3s ease;
    }
    
    .stFileUploader:hover {
        border-color: #667eea;
        background: #f7fafc;
    }
    
    /* Success/Error/Warning messages */
    .stAlert {
        border-radius: 8px;
        border: none;
        padding: 1rem 1.5rem;
    }
    
    /* Chat messages */
    .stChatMessage {
        border-radius: 12px;
        padding: 1rem;
        margin-bottom: 1rem;
        box-shadow: 0 2px 4px rgba(0, 0, 0, 0.05);
    }
    
    /* Animated gradient background */
    .hero-section {
        background: linear-gradient(-45deg, #667eea, #764ba2, #f687b3, #667eea);
        background-size: 400% 400%;
        animation: gradient 15s ease infinite;
        padding: 3rem 2rem;
        border-radius: 16px;
        color: white;
        margin-bottom: 2rem;
        text-align: center;
    }
    
    @keyframes gradient {
        0% { background-position: 0% 50%; }
        50% { background-position: 100% 50%; }
        100% { background-position: 0% 50%; }
    }
    
    /* Loading animation */
    .loading-dots {
        display: inline-block;
        animation: loading 1.4s infinite ease-in-out both;
    }
    
    @keyframes loading {
        0%, 80%, 100% { opacity: 0; }
        40% { opacity: 1; }
    }
    </style>
    """, unsafe_allow_html=True)

# ========== Azure OpenAI Helpers ==========
@st.cache_resource(show_spinner=False)
def get_client():
    if not AZURE_API_KEY or not AZURE_ENDPOINT:
        st.error("🚨 Missing Azure credentials. Set env vars.")
        st.stop()
    return AzureOpenAI(
        api_key        = AZURE_API_KEY,
        azure_endpoint = AZURE_ENDPOINT,
        api_version    = API_VERSION,
    )

def chat(system: str, user: str, temperature: float = 0.3, max_attempts: int = 2) -> str:
    client = get_client()
    for attempt in range(max_attempts):
        try:
            resp = client.chat.completions.create(
                model       = CHAT_MODEL,
                messages    = [
                    {"role": "system", "content": system},
                    {"role": "user",   "content": user},
                ],
                temperature = temperature,
            )
            return resp.choices[0].message.content.strip()
        except Exception as e:
            if attempt + 1 == max_attempts:
                raise e
            time.sleep(1.5)

# ========== File Processing ==========
def extract_guidelines(file):
    if file is None:
        return ""
    try:
        if file.type == "application/pdf":
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() or ""
            return text
        elif file.type in ["text/plain", "text/markdown"]:
            return file.read().decode("utf-8")
        else:
            return ""
    except Exception as e:
        return f"Could not read guidelines: {e}"

def extract_slide_data(ppt_io: BytesIO) -> List[Dict[str, Any]]:
    prs = Presentation(ppt_io)
    data = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                title = sh.text_frame.text.strip()
                break
        text = [sh.text for sh in slide.shapes if hasattr(sh, "text") and sh.text.strip()]
        notes = ""
        if getattr(slide, "notes_slide", None):
            for sh in slide.notes_slide.shapes:
                if hasattr(sh, "text") and sh.text.strip():
                    notes += sh.text + "\n"
        data.append({
            "slide_number": i,
            "title": title,
            "slide_text": "\n".join(text),
            "notes_text": notes.strip(),
            "has_notes": bool(notes.strip()),
        })
    return data

# ========== Helper Functions ==========
def chat_with_guidelines(system, user, guidelines, temperature=0.3, max_attempts=2):
    if guidelines:
        user = f"[CLIENT DESIGN GUIDELINES]\n{guidelines.strip()}\n\n[REQUEST]\n{user}"
    return chat(system, user, temperature, max_attempts)

def extract_json(text):
    match = re.search(r'\{[\s\S]*\}', text)
    if match:
        try:
            return json.loads(match.group(0))
        except Exception:
            return {}
    return {}

def extract_json_array(text):
    """Extract JSON array from text"""
    match = re.search(r'\[[\s\S]*\]', text)
    if match:
        try:
            return json.loads(match.group(0))
        except Exception:
            return []
    return []

# ========== Optimized Analysis Functions ==========
@st.cache_data(show_spinner=False, ttl=3600)
def analyze_deck_batch(slides: List[Dict[str, Any]], guidelines: str = "") -> Dict[str, Any]:
    """Analyze entire deck in a single API call"""
    
    # Prepare the full deck content
    full_content = []
    for d in slides:
        segment = f"[SLIDE {d['slide_number']}]\n"
        segment += f"Title: {d['title'] or '(no title)'}\n"
        segment += f"Content: {d['slide_text']}"
        if d["has_notes"]:
            segment += f"\nNotes: {d['notes_text']}"
        full_content.append(segment)
    
    deck_text = "\n\n".join(full_content)
    
    # Create a comprehensive prompt for all analyses
    prompt = f"""Analyze this presentation and provide a comprehensive JSON response with the following structure:

{{
    "context": {{
        "topic": "main topic of the presentation",
        "region": "target region or market",
        "purpose": "primary purpose of the presentation"
    }},
    "deck_summary": "3 punchy bullet points: 1) Main argument, 2) Target audience, 3) Action items",
    "slides": [
        {{
            "slide_number": 1,
            "summary": "Two sentence summary focusing on essentials",
            "persuasion_tip": "One suggestion for making it more persuasive",
            "layout": "Recommended layout (e.g., 'Title + 2 column layout', 'Full image with text overlay', etc.)",
            "chart_type": "Suggested chart type if applicable (e.g., 'Bar chart', 'Line graph', 'No chart needed')",
            "colors": ["#hex1", "#hex2", "#hex3", "#hex4"] // 4-5 modern hex color codes
        }},
        // ... for each slide
    ]
}}

Presentation content:
{deck_text}"""

    response = chat_with_guidelines(
        "You are an expert presentation analyst. Analyze the presentation and return ONLY valid JSON.",
        prompt,
        guidelines,
        temperature=0.3
    )
    
    # Extract JSON from response
    result = extract_json(response)
    
    # Fallback if parsing fails
    if not result or "slides" not in result:
        st.error("Failed to parse analysis results. Using fallback method.")
        return analyze_deck_fallback(slides, guidelines)
    
    return result

def analyze_deck_fallback(slides: List[Dict[str, Any]], guidelines: str = "") -> Dict[str, Any]:
    """Fallback method with minimal API calls if batch analysis fails"""
    
    # Get context from first few slides (1 API call)
    context = identify_context(slides)
    
    # Get deck summary (1 API call)
    deck_summary = get_deck_summary(slides, guidelines)
    
    # Analyze all slides in batches of 5 (reduces API calls significantly)
    slide_analyses = []
    batch_size = 5
    
    for i in range(0, len(slides), batch_size):
        batch = slides[i:i+batch_size]
        batch_content = []
        
        for d in batch:
            content = f"[SLIDE {d['slide_number']}]\n"
            content += f"Content: {d['slide_text']}"
            if d["has_notes"]:
                content += f"\nNotes: {d['notes_text']}"
            batch_content.append(content)
        
        prompt = f"""For each slide below, provide a JSON array with:
- summary: Two sentence summary
- persuasion_tip: One suggestion for making it more persuasive
- layout: Recommended layout
- chart_type: Suggested chart type if applicable
- colors: Array of 4-5 hex color codes

Slides:
{chr(10).join(batch_content)}

Return ONLY a JSON array of objects."""

        response = chat_with_guidelines(
            "You are a presentation analyst. Return only valid JSON.",
            prompt,
            guidelines,
            temperature=0.3
        )
        
        batch_analyses = extract_json_array(response)
        slide_analyses.extend(batch_analyses)
    
    # Construct the result
    result = {
        "context": context,
        "deck_summary": deck_summary,
        "slides": []
    }
    
    for i, analysis in enumerate(slide_analyses[:len(slides)]):
        result["slides"].append({
            "slide_number": i + 1,
            "summary": analysis.get("summary", ""),
            "persuasion_tip": analysis.get("persuasion_tip", ""),
            "layout": analysis.get("layout", ""),
            "chart_type": analysis.get("chart_type", ""),
            "colors": analysis.get("colors", ["#667eea", "#764ba2", "#f687b3", "#fed7d7", "#2d3748"])
        })
    
    return result

def identify_context(slides: List[Dict[str, Any]]) -> Dict[str, str]:
    """Identify context from first few slides"""
    block = []
    for d in slides[:min(5, len(slides))]:
        b = f"Slide {d['slide_number']} content:\n{d['slide_text']}"
        if d["has_notes"]:
            b += f"\n\nNotes:\n{d['notes_text']}"
        block.append(b)
    
    prompt = (
        "From the following slides, infer the topic, region, and purpose. "
        "Respond ONLY with valid JSON: "
        '{"topic": "...", "region": "...", "purpose": "..."}\n\n'
        "Slides:\n" + "\n\n".join(block)
    )
    
    raw = chat("You are an expert presentation analyst.", prompt, 0.35)
    context = extract_json(raw)
    
    return {
        "topic": context.get("topic", "Unknown"),
        "region": context.get("region", "Unknown"),
        "purpose": context.get("purpose", "Unknown"),
    }

def get_deck_summary(slides: List[Dict[str, Any]], guidelines: str = "") -> str:
    """Get deck summary"""
    full = []
    for d in slides[:10]:  # Limit to first 10 slides for summary
        segment = f"Slide {d['slide_number']}: {d['slide_text'][:200]}..."
        full.append(segment)
    
    return chat_with_guidelines(
        "Summarize the deck in 3 punchy bullet points: 1) Main argument, 2) Target audience, 3) Action items.",
        "\n\n".join(full),
        guidelines,
        0.35
    )

# ========== UI Components ==========
def render_color_palette(colors):
    if not colors:
        colors = ["#667eea", "#764ba2", "#f687b3", "#fed7d7", "#2d3748"]
    
    palette_html = '<div class="color-palette">'
    for code in colors:
        palette_html += f'<div class="color-swatch" style="background-color: {code};"><span class="color-hex">{code}</span></div>'
    palette_html += '</div>'
    
    st.markdown(palette_html, unsafe_allow_html=True)

def render_context_badges(ctx):
    badge_colors = ["", "purple", "blue"]
    badges_html = ""
    for idx, (key, value) in enumerate(ctx.items()):
        color_class = badge_colors[idx % len(badge_colors)]
        badges_html += f'<span class="context-badge {color_class}">{key}: {value}</span>'
    st.markdown(badges_html, unsafe_allow_html=True)

def render_metric_cards(slides):
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.markdown(f"""
        <div class="metric-card">
            <div class="metric-value">{len(slides)}</div>
            <div class="metric-label">Total Slides</div>
        </div>
        """, unsafe_allow_html=True)
    
    with col2:
        notes_count = sum(s['has_notes'] for s in slides)
        st.markdown(f"""
        <div class="metric-card" style="background: linear-gradient(135deg, #48bb78 0%, #38a169 100%);">
            <div class="metric-value">{notes_count}</div>
            <div class="metric-label">With Notes</div>
        </div>
        """, unsafe_allow_html=True)
    
    with col3:
        avg_words = sum(len(s['slide_text'].split()) for s in slides) // len(slides) if slides else 0
        st.markdown(f"""
        <div class="metric-card" style="background: linear-gradient(135deg, #ed8936 0%, #dd6b20 100%);">
            <div class="metric-value">{avg_words}</div>
            <div class="metric-label">Avg Words/Slide</div>
        </div>
        """, unsafe_allow_html=True)
    
    with col4:
        est_time = max(1, len(slides) // 2)
        st.markdown(f"""
        <div class="metric-card" style="background: linear-gradient(135deg, #9f7aea 0%, #805ad5 100%);">
            <div class="metric-value">{est_time}</div>
            <div class="metric-label">Est. Minutes</div>
        </div>
        """, unsafe_allow_html=True)

# ========== AI Sanitization ==========
def sanitize_content_batch(texts: List[str]) -> List[str]:
    """Sanitize multiple texts in a single API call"""
    if not texts or all(not t.strip() for t in texts):
        return texts
    
    # Create a batch prompt
    batch_prompt = "Sanitize each of the following texts by removing sensitive content (names, companies, dates, financial figures). " \
                   "Replace with placeholders like {Name}, {Company}, {Date}. " \
                   "Maintain original tone and structure. Return a JSON array of sanitized texts.\n\n"
    
    for i, text in enumerate(texts):
        batch_prompt += f"[TEXT {i+1}]\n{text}\n\n"
    
    batch_prompt += "Return ONLY a JSON array of sanitized texts in the same order."
    
    response = chat(
        "You are a content anonymization assistant.",
        batch_prompt,
        temperature=0.2
    )
    
    sanitized = extract_json_array(response)
    
    # Fallback to individual sanitization if batch fails
    if not sanitized or len(sanitized) != len(texts):
        return [sanitize_content_with_ai(t) for t in texts]
    
    return sanitized

def sanitize_content_with_ai(text):
    """Fallback individual sanitization"""
    if not text or not text.strip():
        return text

    word_count = len(text.split())
    max_allowed_words = int(word_count * 1.2)

    prompt = (
        f"Sanitize the following text by removing only actual sensitive content. "
        f"Replace sensitive info with placeholders like {{Name}}, {{Company}}, {{Date}}. "
        f"Keep under {max_allowed_words} words.\n\n{text}"
    )

    response = chat(
        "You are a content anonymization assistant.",
        prompt,
        temperature=0.2
    )

    sanitized_words = response.split()
    if len(sanitized_words) > max_allowed_words:
        response = " ".join(sanitized_words[:max_allowed_words]) + "..."

    return response

def sanitize_pptx_optimized(ppt_file: BytesIO) -> BytesIO:
    """Optimized sanitization with batch processing and progress tracking"""
    prs = Presentation(ppt_file)
    
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    # Collect all texts to sanitize
    texts_to_sanitize = []
    text_locations = []  # Keep track of where each text came from
    
    status_text.text("Collecting texts from presentation...")
    
    for slide_idx, slide in enumerate(prs.slides):
        # Collect slide texts
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        texts_to_sanitize.append(text)
                        text_locations.append(("slide", slide_idx, shape_idx, para_idx))
        
        # Collect notes texts
        if getattr(slide, "notes_slide", None):
            for shape_idx, shape in enumerate(slide.notes_slide.shapes):
                if shape.has_text_frame:
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            texts_to_sanitize.append(text)
                            text_locations.append(("notes", slide_idx, shape_idx, para_idx))
        
        progress_bar.progress((slide_idx + 1) / len(prs.slides) * 0.3)
    
    # Sanitize all texts in batches
    batch_size = 10
    sanitized_texts = []
    
    status_text.text("Sanitizing content...")
    
    for i in range(0, len(texts_to_sanitize), batch_size):
        batch = texts_to_sanitize[i:i+batch_size]
        sanitized_batch = sanitize_content_batch(batch)
        sanitized_texts.extend(sanitized_batch)
        
        progress = 0.3 + (0.5 * ((i + batch_size) / len(texts_to_sanitize)))
        progress_bar.progress(min(progress, 0.8))
        status_text.text(f"Sanitizing content... {min(i + batch_size, len(texts_to_sanitize))}/{len(texts_to_sanitize)} items")
    
    # Apply sanitized texts back to presentation
    status_text.text("Applying sanitized content...")
    
    for i, (location_type, slide_idx, shape_idx, para_idx) in enumerate(text_locations):
        sanitized_text = sanitized_texts[i]
        
        if location_type == "slide":
            slide = prs.slides[slide_idx]
            shape = slide.shapes[shape_idx]
            para = shape.text_frame.paragraphs[para_idx]
        else:  # notes
            slide = prs.slides[slide_idx]
            shape = slide.notes_slide.shapes[shape_idx]
            para = shape.text_frame.paragraphs[para_idx]
        
        # Preserve formatting
        bullet = para.level
        align = para.alignment
        first_run = para.runs[0] if para.runs else None
        
        para.clear()
        new_run = para.add_run()
        new_run.text = sanitized_text
        
        if first_run:
            copy_font_style(first_run, new_run)
        para.level = bullet
        para.alignment = align
        
        progress = 0.8 + (0.2 * ((i + 1) / len(text_locations)))
        progress_bar.progress(progress)
    
    progress_bar.empty()
    status_text.empty()
    
    sanitized_io = BytesIO()
    prs.save(sanitized_io)
    sanitized_io.seek(0)
    return sanitized_io

def copy_font_style(source_run, target_run):
    if source_run.font is None:
        return
    target_font = target_run.font
    source_font = source_run.font

    target_font.name = source_font.name
    target_font.size = source_font.size
    target_font.bold = source_font.bold
    target_font.italic = source_font.italic

    try:
        if source_font.color and source_font.color.rgb:
            rgb = source_font.color.rgb
            target_font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
    except AttributeError:
        pass

def generate_visual_prompt(slide_summary: str, context: Dict[str, str], guidelines: str = "") -> str:
    system = (
        "Act as a design consultant for executive presentations. Suggest a creative, non-generic image prompt "
        "for DALL-E, making it region- and topic-specific. Return just the image prompt."
    )
    user = (
        f"Create an image prompt for this slide:\n\n{slide_summary}\n\n"
        f"Context: {json.dumps(context)}"
    )
    return chat_with_guidelines(system, user, guidelines, 0.7)

def build_chat_corpus(analysis: Dict[str, Any], slides: List[Dict[str, Any]]) -> str:
    """Build chat corpus from analysis results"""
    bits = [
        f"Presentation context: {json.dumps(analysis['context'])}",
        f"Deck summary: {analysis['deck_summary']}"
    ]
    
    for slide_data, slide_analysis in zip(slides, analysis['slides']):
        block = [
            f"Slide {slide_data['slide_number']}",
            f"Title: {slide_data['title']}",
            f"Content: {slide_data['slide_text']}"
        ]
        if slide_data["has_notes"]:
            block.append(f"Notes: {slide_data['notes_text']}")
        
        block.extend([
            f"AI Summary: {slide_analysis['summary']}",
            f"Design tips: {slide_analysis['persuasion_tip']}",
            f"Layout suggestion: {slide_analysis['layout']}",
            f"Chart suggestion: {slide_analysis['chart_type']}",
        ])
        bits.append("\n".join(block))
    
    return "\n\n".join(bits)

# ========== Main Application ==========
def main():
    st.set_page_config(
        page_title="PPT Analyzer Pro", 
        layout="wide", 
        page_icon="🎯",
        initial_sidebar_state="expanded"
    )
    
    # Inject custom CSS
    inject_custom_css()
    
    # Hero Section
    st.markdown("""
    <div class="hero-section">
        <h1 style="font-size: 3.5rem; margin-bottom: 1rem;">PPT Analyzer Pro</h1>
        <p style="font-size: 1.25rem; opacity: 0.95;">Transform your presentations with AI-powered insights</p>
    </div>
    """, unsafe_allow_html=True)
    
    # Sidebar
    with st.sidebar:
        st.markdown("## 📁 Upload Files")
        
        ppt = st.file_uploader(
            "Choose PowerPoint", 
            type=["pptx"],
            help="Upload your presentation file"
        )
        
        guidelines_file = st.file_uploader(
            "Design Guidelines (Optional)", 
            type=["pdf", "txt", "md"],
            help="Upload brand or design guidelines"
        )
        
        st.markdown("---")
        
        # Azure Status
        st.markdown("### 🔌 Connection Status")
        if AZURE_API_KEY and AZURE_ENDPOINT:
            st.success("✅ Azure OpenAI Connected")
        else:
            st.error("❌ Azure OpenAI Not Connected")
            st.info("Please set your Azure OpenAI credentials in the environment variables")
        
        st.markdown("---")
        
        # Action Buttons
        st.markdown("### 🎯 Actions")
        col1, col2 = st.columns(2)
        with col1:
            analyze = st.button("🔍 Analyze", use_container_width=True)
        with col2:
            sanitize = st.button("🧹 Sanitize", use_container_width=True)
    
    # Main content area
    if ppt is None:
        # Welcome screen
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.markdown("""
            <div class="deck-card">
                <h3>📊 Smart Analysis</h3>
                <p>Get AI-powered insights on your slides, including summaries, design tips, and visual recommendations.</p>
            </div>
            """, unsafe_allow_html=True)
        
        with col2:
            st.markdown("""
            <div class="deck-card">
                <h3>🎨 Design Assistant</h3>
                <p>Receive layout suggestions, color palettes, and chart recommendations tailored to each slide.</p>
            </div>
            """, unsafe_allow_html=True)
        
        with col3:
            st.markdown("""
            <div class="deck-card">
                <h3>💬 Interactive Chat</h3>
                <p>Ask questions about your presentation and get instant, context-aware answers.</p>
            </div>
            """, unsafe_allow_html=True)
        
        st.info("👆 Upload a PowerPoint file in the sidebar to get started")
        return
    
    # Process guidelines
    guidelines = extract_guidelines(guidelines_file) if guidelines_file else ""
    
    # Session state management
    if "ppt_name" not in st.session_state:
        st.session_state.ppt_name = None
    if "analysis_ready" not in st.session_state:
        st.session_state.analysis_ready = False
    
    # Check if file changed
    if ppt.name != st.session_state.ppt_name:
        st.session_state.analysis_ready = False
        st.session_state.ppt_name = ppt.name
    
    # Handle sanitization
    if sanitize:
        with st.spinner("🧹 Sanitizing your presentation..."):
            try:
                sanitized_io = sanitize_pptx_optimized(ppt)
                st.success("✨ Your presentation has been sanitized!")
                st.download_button(
                    label="📥 Download Sanitized PPT",
                    data=sanitized_io,
                    file_name="Sanitized_Presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
            except Exception as e:
                st.error(f"❌ Error during sanitization: {str(e)}")
        return
    
    # Handle analysis
    if analyze and ppt is not None:
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        try:
            # Extract slides
            status_text.text("📖 Reading presentation...")
            progress_bar.progress(0.1)
            slides = extract_slide_data(ppt)
            
            # Analyze deck
            status_text.text("🧠 Analyzing presentation content...")
            progress_bar.progress(0.3)
            
            with st.spinner("This may take a moment for large presentations..."):
                analysis = analyze_deck_batch(slides, guidelines)
            
            progress_bar.progress(0.9)
            status_text.text("✅ Finalizing analysis...")
            
            # Store results
            st.session_state.analysis = analysis
            st.session_state.slides = slides
            st.session_state.corpus = build_chat_corpus(analysis, slides)
            st.session_state.analysis_ready = True
            
            progress_bar.progress(1.0)
            time.sleep(0.5)
            progress_bar.empty()
            status_text.empty()
            
            st.success("✅ Analysis complete!")
            
        except Exception as e:
            progress_bar.empty()
            status_text.empty()
            st.error(f"❌ Error during analysis: {str(e)}")
            st.info("Please try again or check your Azure OpenAI connection.")
            return
    
    if not st.session_state.get("analysis_ready", False):
        st.warning("👈 Click 'Analyze' in the sidebar to process your presentation")
        return
    
    # Display results
    slides = st.session_state.slides
    analysis = st.session_state.analysis
    
    # Create tabs
    tab1, tab2, tab3 = st.tabs(["📊 Overview", "🎯 Slide Details", "💬 AI Assistant"])
    
    with tab1:
        # Metrics
        render_metric_cards(slides)
        
        # Context and Summary
        st.markdown("<div class='deck-card'>", unsafe_allow_html=True)
        st.markdown("### 🎯 Presentation Context")
        render_context_badges(analysis['context'])
        
        st.markdown("### 📝 Executive Summary")
        st.markdown(f'<div class="summary-box">{analysis["deck_summary"]}</div>', unsafe_allow_html=True)
        st.markdown("</div>", unsafe_allow_html=True)
    
    with tab2:
        # Slide selector
        selected = st.selectbox(
            "Select a slide to analyze:",
            range(len(slides)),
            format_func=lambda x: f"Slide {x+1}: {slides[x]['title'][:50] or '(Untitled)'}...",
            key="slide_selector"
        )
        
        slide = slides[selected]
        slide_analysis = analysis["slides"][selected]
        
        # Slide content in card
        st.markdown("<div class='deck-card'>", unsafe_allow_html=True)
        
        col1, col2 = st.columns([2, 1])
        
        with col1:
            st.markdown(f"### Slide {slide['slide_number']}: {slide['title'] or '(Untitled)'}")
            
            st.markdown("**Content:**")
            st.info(slide['slide_text'] or "No text content")
            
            if slide['has_notes']:
                with st.expander("📝 Speaker Notes"):
                    st.text(slide['notes_text'])
        
        with col2:
            st.markdown("### 🎨 Design Elements")
            
            st.markdown("**Layout:**")
            st.success(slide_analysis['layout'])
            
            st.markdown("**Visualization:**")
            st.info(slide_analysis['chart_type'])
            
            st.markdown("**Color Palette:**")
            render_color_palette(slide_analysis['colors'])
        
        st.markdown("</div>", unsafe_allow_html=True)
        
        # AI Insights
        st.markdown("<div class='deck-card'>", unsafe_allow_html=True)
        st.markdown("### 🤖 AI Analysis")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("**Summary:**")
            st.markdown(f'<div class="summary-box">{slide_analysis["summary"]}</div>', unsafe_allow_html=True)
        
        with col2:
            st.markdown("**Enhancement Suggestion:**")
            st.markdown(f'<div class="ai-suggestion">{slide_analysis["persuasion_tip"]}</div>', unsafe_allow_html=True)
        
        # Visual prompt generator
        if st.button("🎨 Generate Visual Concept", key="visual_btn"):
            with st.spinner("Creating visual concept..."):
                prompt = generate_visual_prompt(
                    f"{slide_analysis['summary']} {slide_analysis['persuasion_tip']}", 
                    analysis['context'], 
                    guidelines
                )
                st.markdown("**🖼️ DALL-E Prompt:**")
                st.code(prompt, language="text")
        
        st.markdown("</div>", unsafe_allow_html=True)
    
    with tab3:
        st.markdown("### 💬 Presentation Assistant")
        st.caption("Ask me anything about your presentation!")
        
        # Chat interface
        if "messages" not in st.session_state:
            st.session_state.messages = []
        
        # Display chat history
        for message in st.session_state.messages:
            with st.chat_message(message["role"]):
                st.write(message["content"])
        
        # Chat input
        if prompt := st.chat_input("What would you like to know about your presentation?"):
            # Add user message
            st.session_state.messages.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.write(prompt)
            
            # Generate response
            with st.chat_message("assistant"):
                with st.spinner("Thinking..."):
                    response = chat_with_guidelines(
                        system="You are a helpful presentation assistant. Answer based on the presentation content and guidelines provided.",
                        user=f"{prompt}\n\nPresentation content:\n{st.session_state.corpus}",
                        guidelines=guidelines
                    )
                    st.write(response)
            
            st.session_state.messages.append({"role": "assistant", "content": response})

if __name__ == "__main__":
    main()